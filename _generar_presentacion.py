from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ANTRACITA = RGBColor(0x2C, 0x35, 0x39)
TERRACOTA = RGBColor(0xE6, 0x4A, 0x19)
MARFIL = RGBColor(0xF7, 0xF8, 0xF7)
GRIS = RGBColor(0x5A, 0x66, 0x6B)
FONDO = RGBColor(0xEC, 0xEF, 0xF1)
ROJO = RGBColor(0xC6, 0x28, 0x28)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
FUENTE = "Segoe UI"

prs = Presentation()
prs.slide_width = 12192000
prs.slide_height = 6858000
BLANK = prs.slide_layouts[6]


def nueva():
    s = prs.slides.add_slide(BLANK)
    fondo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = FONDO
    fondo.line.fill.background()
    fondo.shadow.inherit = False
    return s


def barra_titulo(s, titulo):
    banda = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, 990600)
    banda.fill.solid()
    banda.fill.fore_color.rgb = TERRACOTA
    banda.line.fill.background()
    banda.shadow.inherit = False
    caja = s.shapes.add_textbox(450000, 180000, 11000000, 700000)
    p = caja.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = titulo
    r.font.size = 300000
    r.font.bold = True
    r.font.name = FUENTE
    r.font.color.rgb = BLANCO


def texto(s, x, y, ancho, alto, lineas, tam=16, color=ANTRACITA, negrita=False,
          alineacion=PP_ALIGN.LEFT, interlineado=1.15):
    caja = s.shapes.add_textbox(x, y, ancho, alto)
    tf = caja.text_frame
    tf.word_wrap = True
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alineacion
        p.line_spacing = interlineado
        if isinstance(linea, tuple):
            contenido, es_negrita = linea
        else:
            contenido, es_negrita = linea, negrita
        r = p.add_run()
        r.text = contenido
        r.font.size = tam * 12700
        r.font.name = FUENTE
        r.font.bold = es_negrita
        r.font.color.rgb = color
    return caja


def tarjeta(s, x, y, ancho, alto, titulo, cuerpo, color_titulo=TERRACOTA):
    caja = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, ancho, alto)
    caja.adjustments[0] = 0.06
    caja.fill.solid()
    caja.fill.fore_color.rgb = MARFIL
    caja.line.color.rgb = RGBColor(0xDD, 0xE2, 0xE5)
    caja.shadow.inherit = False
    tf = caja.text_frame
    tf.word_wrap = True
    tf.margin_left = 200000
    tf.margin_top = 140000
    pt = tf.paragraphs[0]
    rt = pt.add_run()
    rt.text = titulo
    rt.font.size = 170000
    rt.font.bold = True
    rt.font.name = FUENTE
    rt.font.color.rgb = color_titulo
    for linea in cuerpo:
        p = tf.add_paragraph()
        p.space_before = 50000
        r = p.add_run()
        r.text = linea
        r.font.size = 130000
        r.font.name = FUENTE
        r.font.color.rgb = ANTRACITA
    return caja


def caja_flujo(s, x, y, ancho, alto, titulo, subtitulo):
    forma = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, ancho, alto)
    forma.adjustments[0] = 0.10
    forma.fill.solid()
    forma.fill.fore_color.rgb = ANTRACITA
    forma.line.fill.background()
    forma.shadow.inherit = False
    tf = forma.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = titulo
    r1.font.size = 160000
    r1.font.bold = True
    r1.font.name = FUENTE
    r1.font.color.rgb = BLANCO
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = 40000
    r2 = p2.add_run()
    r2.text = subtitulo
    r2.font.size = 120000
    r2.font.name = FUENTE
    r2.font.color.rgb = RGBColor(0xB8, 0xC4, 0xC9)


def flecha(s, x, y, ancho=450000):
    f = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, ancho, 350000)
    f.fill.solid()
    f.fill.fore_color.rgb = TERRACOTA
    f.line.fill.background()
    f.shadow.inherit = False


def tabla(s, filas, x, y, ancho, alto):
    cols = len(filas[0])
    shape = s.shapes.add_table(len(filas), cols, x, y, ancho, alto)
    t = shape.table
    for i, fila in enumerate(filas):
        for j, celda in enumerate(fila):
            cell = t.cell(i, j)
            cell.text = str(celda)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = 100000
            par = cell.text_frame.paragraphs[0]
            run = par.runs[0]
            run.font.name = FUENTE
            run.font.size = 120000
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TERRACOTA
                run.font.color.rgb = BLANCO
                run.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = MARFIL if i % 2 else BLANCO
                run.font.color.rgb = ANTRACITA
    return t


s = nueva()
texto(s, 900000, 2100000, 10300000, 1500000, [("NUAM", True)], tam=54, color=TERRACOTA)
texto(s, 900000, 3300000, 10300000, 900000, [("Mantenedor de Calificaciones Tributarias", True)], tam=30)
texto(s, 900000, 4200000, 10300000, 700000,
      ["Sistema web en Flask + MongoDB para la homologación DJ 1949 → DJ 1922",
       "Conversión automática de montos a factores tributarios"], tam=16, color=GRIS)
linea = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 950000, 3150000, 2400000, 60000)
linea.fill.solid()
linea.fill.fore_color.rgb = TERRACOTA
linea.line.fill.background()
linea.shadow.inherit = False
texto(s, 900000, 5800000, 10300000, 500000,
      ["github.com/Pipstachio77777dev/NUAM-ANTO"], tam=13, color=GRIS)

s = nueva()
barra_titulo(s, "Contexto y problema")
tarjeta(s, 500000, 1500000, 5400000, 2200000, "El negocio",
        ["Los corredores de bolsa deben calificar tributariamente",
         "los dividendos pagados por cada instrumento en un ejercicio.",
         "",
         "La homologación DJ 1949 → DJ 1922 se resuelve con 30 columnas",
         "codificadas (8 a 37), cada una con un factor tributario."])
tarjeta(s, 6200000, 1500000, 5400000, 2200000, "El dolor actual",
        ["Cálculo manual propenso a errores de redondeo.",
         "Sin historial: no se sabe quién creó o modificó qué.",
         "Cargas masivas en planillas dispersas y sin validación."])
tarjeta(s, 500000, 3950000, 11100000, 1900000, "La solución NUAM",
        ["Mantenedor web centralizado que calcula los factores con reglas exactas e inmutables,",
         "valida los datos al ingresarlos, acepta cargas masivas por CSV",
         "y deja trazabilidad completa de cada acción (crear / modificar / eliminar / carga)."])

s = nueva()
barra_titulo(s, "Arquitectura técnica")
caja_flujo(s, 700000, 1900000, 2900000, 1300000, "Navegador", "SPA HTML/CSS/JS sin frameworks")
flecha(s, 3750000, 2370000)
caja_flujo(s, 4350000, 1900000, 3200000, 1300000, "API REST Flask", "Python 3.14 · venv · puerto 5000")
flecha(s, 7700000, 2370000)
caja_flujo(s, 8300000, 1900000, 3100000, 1300000, "MongoDB", "Colecciones calificaciones y auditoria")
texto(s, 700000, 3600000, 10800000, 2500000,
      [("", False),
       ("• dominio/ — lógica tributaria pura (factores.py) y validadores con mensajes exactos", False),
       ("• db.py — índice único (ejercicio + mercado + instrumento + secuencia), semilla automática", False),
       ("• app.py — API REST completa con auditoría automática de cada operación", False),
       ("• static/ y templates/ — interfaz SPA con navegación por estado y fetch", False),
       ("• Scripts .bat — requerimientos.bat (instala) e iniciar.bat (levanta y abre el navegador)", False)],
      tam=15)

s = nueva()
barra_titulo(s, "Reglas del dominio tributario")
tarjeta(s, 500000, 1450000, 5450000, 2450000, "Cálculo del factor (columna i)",
        ["peso_i = 0.96^i  (i = 0..29), normalizado y redondeado a 8 decimales",
         "",
         "actualización = factor_actualizacion si > 0, si no 1",
         "valorActualizado = valor_historico × actualización",
         "",
         "factor = round8( (dividendo × peso_i) ÷ valorActualizado )",
         "clamp SIEMPRE a [0, 1]"])
tarjeta(s, 6250000, 1450000, 5450000, 2450000, "Redondeo round8 sin errores flotantes",
        ["Formatea a 9 decimales, observa el dígito 9°:",
         "si ≥ 5 suma 1 al bloque de 8 decimales (con acarreo).",
         "",
         "round8(0.123456789) → 0.12345679",
         "round8(0.1 + 0.2)   → 0.3"])
tarjeta(s, 500000, 4150000, 11200000, 1650000, "Valor verificado en pruebas",
        ["EMPRESAS CMPC S.A. (dividendo 1.200,5 · valor histórico 12.500 · factor actualización 1,00024)",
         "→ factor columna 9 = 0.0052214  ✓  (coincide con el criterio de aceptación del negocio)",
         "Los factores se persisten como JSON {\"8\": …, \"37\": …} y se recalculan al guardar."])

s = nueva()
barra_titulo(s, "Pantallas del sistema")
datos_pantallas = [
    ("Inicio", "Acceso rápido a las 3 funciones principales."),
    ("Consultar", "Filtros por mercado (ACC/CFI/FMU), origen y ejercicio; grilla de resultados."),
    ("Ingresar / Modificar", "Formulario completo; muestra los errores exactos de la API."),
    ("Carga masiva", "Pega líneas o sube un CSV; resumen insertados vs. con errores por fila."),
    ("Trazabilidad", "Bitácora con código de color: terracota = acciones, rojo = eliminaciones."),
]
y = 1350000
for titulo, descripcion in datos_pantallas:
    tarjeta(s, 500000, y, 11200000, 850000, titulo, [descripcion])
    y += 1000000
texto(s, 500000, y + 50000, 11200000, 500000,
      ["Login local fijo: ADMIN@mail.com / 1234 · Sidebar colapsable que nunca desplaza el contenido"],
      tam=13, color=GRIS)

s = nueva()
barra_titulo(s, "API REST")
tabla(s, [
    ["Método", "Ruta", "Descripción"],
    ["GET", "/api/calificaciones?mercado=&origen=&ejercicio=", "Lista con filtros opcionales"],
    ["POST", "/api/calificaciones", "Crear → 201 · errores → 400 {\"errores\":[…]}"],
    ["PUT", "/api/calificaciones/<id>", "Modificar → 200 · 404 si no existe"],
    ["DELETE", "/api/calificaciones/<id>", "Eliminar → 204 vacío · 404 si no existe"],
    ["GET", "/api/calificaciones/catalogos", "Ejercicios sin duplicados: [2022, 2023, 2024]"],
    ["GET", "/api/calificaciones/auditoria", "Bitácora, más recientes primero"],
    ["POST", "/api/calificaciones/masiva", "Carga masiva {\"lineas\":[…]} con detalle por fila"],
], 500000, 1500000, 11200000, 3800000)
texto(s, 500000, 5700000, 11200000, 700000,
      ["Formato camelCase · auditoría automática con usuario y resumen por evento"],
      tam=14, color=GRIS)

s = nueva()
barra_titulo(s, "Base de datos MongoDB")
tarjeta(s, 500000, 1450000, 5450000, 2600000, "colección: calificaciones",
        ["mercado (ACC/CFI/FMU) · instrumento · secuencia_evento",
         "dividendo · fecha_pago DD-MM-AAAA · valor_historico",
         "ejercicio · isfut · factor_actualizacion · origen (ENTIDAD/SISTEMA)",
         "factores JSON columnas 8–37 · fechas de creación y actualización",
         "",
         "Índice ÚNICO: ejercicio + mercado + instrumento + secuencia_evento"])
tarjeta(s, 6250000, 1450000, 5450000, 2600000, "colección: auditoria",
        ["evento: crear | modificar | eliminar | carga masiva",
         "registro_id · instrumento · ejercicio",
         "resumen legible por humano",
         "usuario (ADMIN@mail.com) · fecha ISO 8601"])
tarjeta(s, 500000, 4300000, 11200000, 1500000, "Semilla automática al primer arranque",
        ["5 registros reales de ejemplo: CMPC, SQM-B, CFI LarrainVial, FM Consultora y FM Principal Chile",
         "+ entrada de auditoría 'Se sembraron 5 registros iniciales de ejemplo'"])

s = nueva()
barra_titulo(s, "Calidad y pruebas automatizadas")
texto(s, 550000, 1400000, 11000000, 800000,
      [("36 verificaciones automáticas sobre los criterios de aceptación — todas en verde", True)],
      tam=19)
items = [
    ("✓", "Factor CMPC columna 9 = 0.0052214 exacto"),
    ("✓", "round8 sin artefactos: 0.12345679 y 0.3 verificados"),
    ("✓", "Todos los factores dentro de [0, 1], columnas 8–37 completas"),
    ("✓", "Catálogos devuelven [2022, 2023, 2024] sin duplicados"),
    ("✓", "Mensajes de validación EXACTOS (ej.: 'Ejercicio debe ser numérico de largo 4')"),
    ("✓", "CRUD completo registra auditoría con resúmenes correctos"),
    ("✓", "Duplicados rechazados (índice único) · PUT/DELETE inexistentes → 404"),
    ("✓", "Carga masiva: línea inválida rechazada, válida insertada con factores"),
]
y = 2300000
for marca, linea_txt in items:
    texto(s, 600000, y, 500000, 420000, [(marca, True)], tam=15, color=TERRACOTA, negrita=True)
    texto(s, 1150000, y, 10500000, 420000, [(linea_txt, False)], tam=15)
    y += 520000

s = nueva()
barra_titulo(s, "Cómo ejecutarlo")
pasos = [
    ("1", "Instalar", "Doble clic en requerimientos.bat (crea el venv e instala dependencias)."),
    ("2", "Requisitos", "MongoDB corriendo en localhost:27017."),
    ("3", "Iniciar", "Doble clic en iniciar.bat → abre http://127.0.0.1:5000 solo."),
    ("4", "Sesión", "Correo ADMIN@mail.com · clave 1234."),
]
y = 1500000
for num, titulo, desc in pasos:
    circulo = s.shapes.add_shape(MSO_SHAPE.OVAL, 650000, y, 620000, 620000)
    circulo.fill.solid()
    circulo.fill.fore_color.rgb = TERRACOTA
    circulo.line.fill.background()
    circulo.shadow.inherit = False
    tfc = circulo.text_frame
    pc = tfc.paragraphs[0]
    pc.alignment = PP_ALIGN.CENTER
    rc = pc.add_run()
    rc.text = num
    rc.font.size = 170000
    rc.font.bold = True
    rc.font.name = FUENTE
    rc.font.color.rgb = BLANCO
    texto(s, 1500000, y - 30000, 10200000, 800000, [(titulo, True), (desc, False)],
          tam=16, color=ANTRACITA)
    y += 1050000
texto(s, 650000, 5900000, 11000000, 500000,
      ["Código fuente: github.com/Pipstachio77777dev/NUAM-ANTO"], tam=13, color=GRIS)

s = nueva()
barra_titulo(s, "Próximos pasos sugeridos")
items_futuro = [
    ("Autenticación real", "Reemplazar el login fijo por usuarios con roles (administrador, operador, consulta)."),
    ("Exportaciones", "Descarga de grillas y factores a Excel/PDF para reportes regulatorios."),
    ("Producción", "Migrar de SQLite-like local a MongoDB Atlas o servidor dedicado con respaldos."),
    ("Más mercados e instrumentos", "Ampliar catálogos según nuevos convenios NUAM (Chile · Colombia · Perú)."),
]
y = 1500000
for titulo, desc in items_futuro:
    tarjeta(s, 500000, y, 11200000, 1000000, titulo, [desc])
    y += 1150000
texto(s, 500000, 6200000, 11200000, 500000,
      ["NUAM — Mantenedor de Calificaciones Tributarias · Gracias"], tam=14, color=GRIS,
      alineacion=PP_ALIGN.CENTER)

prs.save("Presentacion_NUAM.pptx")
print("Presentacion_NUAM.pptx generada con", len(prs.slides.__iter__.__self__._sldIdLst), "diapositivas")
